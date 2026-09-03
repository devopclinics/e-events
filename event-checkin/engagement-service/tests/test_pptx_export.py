"""Tests for the PPTX exporter's pure assembly logic. Deliberately excludes
the browser/screenshot half (_capture_steps) -- that needs a real Chromium
and is covered by manual staging verification instead, matching how this
codebase already treats other browser/network-dependent code."""
import asyncio
import io
import unittest
from unittest.mock import AsyncMock, patch

from fastapi import HTTPException
from pptx import Presentation

from app.pptx_export import SLIDE_HEIGHT_IN, SLIDE_WIDTH_IN, _assemble, _capture_steps, _slide_notes

# The smallest possible valid PNG (1x1, transparent) -- real enough bytes for
# python-pptx's add_picture to accept without needing a rendered screenshot.
_TINY_PNG = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000a49444154789c6360000002000100dd8d1cb0000000004945454e42"
    "601a91821b2c6c0000000049454e44ae426082"
)


class SlideNotesTests(unittest.TestCase):
    def test_includes_the_slide_position_type_and_title(self):
        step = type("Step", (), {
            "sequence": 2, "step_type": "poll_results", "title": "How connected are we?",
            "presenter_notes": "Reveal · pause",
        })()
        notes = _slide_notes(step, index=2, total=27)
        self.assertIn("Slide 3 of 27", notes)
        self.assertIn("poll_results", notes)
        self.assertIn("How connected are we?", notes)
        self.assertIn("Reveal · pause", notes)

    def test_still_labels_the_slide_when_a_step_has_no_presenter_notes(self):
        step = type("Step", (), {"sequence": 0, "step_type": "hero", "title": "Opening", "presenter_notes": None})()
        notes = _slide_notes(step, index=0, total=1)
        self.assertEqual(notes, 'Slide 1 of 1 · hero · "Opening"')


class AssembleDeckTests(unittest.TestCase):
    def test_produces_one_full_bleed_widescreen_slide_per_step_with_matching_notes(self):
        steps = [
            type("Step", (), {"sequence": 0, "step_type": "hero", "title": "THE GOOD LIFE", "presenter_notes": "Assalamu Alaikum..."})(),
            type("Step", (), {"sequence": 1, "step_type": "closing", "title": "BETTER TOGETHER", "presenter_notes": None})(),
        ]

        deck_bytes = _assemble(steps, [_TINY_PNG, _TINY_PNG])

        presentation = Presentation(io.BytesIO(deck_bytes))
        slides = list(presentation.slides)
        self.assertEqual(len(slides), 2)
        self.assertAlmostEqual(presentation.slide_width.inches, SLIDE_WIDTH_IN, places=2)
        self.assertAlmostEqual(presentation.slide_height.inches, SLIDE_HEIGHT_IN, places=2)
        self.assertIn("THE GOOD LIFE", slides[0].notes_slide.notes_text_frame.text)
        self.assertIn("BETTER TOGETHER", slides[1].notes_slide.notes_text_frame.text)

    def test_each_slides_picture_fills_the_entire_frame(self):
        steps = [type("Step", (), {"sequence": 0, "step_type": "hero", "title": "Opening", "presenter_notes": None})()]
        presentation = Presentation(io.BytesIO(_assemble(steps, [_TINY_PNG])))
        picture = list(presentation.slides)[0].shapes[0]
        self.assertEqual(picture.left, 0)
        self.assertEqual(picture.top, 0)
        self.assertEqual(picture.width, presentation.slide_width)
        self.assertEqual(picture.height, presentation.slide_height)


class ConcurrencyGuardTests(unittest.IsolatedAsyncioTestCase):
    """A headless Chromium runs in the same pod as live voting traffic, so a
    second export must be rejected immediately -- not queued behind the
    first for however long that takes, and not left free to stack a second
    browser process alongside it."""

    async def test_a_second_concurrent_export_is_rejected_immediately_not_queued(self):
        release = asyncio.Event()

        async def slow_capture(*_args, **_kwargs):
            await release.wait()
            return [b"fake-screenshot"]

        with patch("app.pptx_export._capture_steps_unguarded", new=AsyncMock(side_effect=slow_capture)):
            first = asyncio.create_task(_capture_steps("http://proxy", "wf-a", "token", []))
            await asyncio.sleep(0.05)  # let the first call claim the export slot
            with self.assertRaises(HTTPException) as raised:
                await _capture_steps("http://proxy", "wf-a", "token", [])
            self.assertEqual(raised.exception.status_code, 503)

            release.set()
            self.assertEqual(await first, [b"fake-screenshot"])

    async def test_the_slot_is_released_after_a_failed_capture_so_the_next_export_can_proceed(self):
        with patch("app.pptx_export._capture_steps_unguarded", new=AsyncMock(side_effect=RuntimeError("boom"))):
            with self.assertRaises(RuntimeError):
                await _capture_steps("http://proxy", "wf-a", "token", [])

        with patch("app.pptx_export._capture_steps_unguarded", new=AsyncMock(return_value=[b"ok"])):
            self.assertEqual(await _capture_steps("http://proxy", "wf-a", "token", []), [b"ok"])


if __name__ == "__main__":
    unittest.main()
