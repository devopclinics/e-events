"""Render a workflow's steps to a PowerPoint deck.

Each slide is a real screenshot of the same WorkflowSceneRenderer the TV
uses -- taken by a headless browser navigating the frontend's step-preview
route (`/live/step-preview/{workflow_id}/{step_id}`) -- so the deck never
drifts from what a presenter actually sees. That route reads a step's
current data straight from the database; it never creates a WorkflowRun or
LiveDisplay, so generating a deck can never interfere with (or be
interfered with by) an active presenter session.
"""
from __future__ import annotations

import asyncio
import io

from fastapi import HTTPException
from playwright.async_api import async_playwright
from pptx import Presentation
from pptx.util import Inches

from .models import WorkflowStep

SLIDE_WIDTH_IN = 13.333
SLIDE_HEIGHT_IN = 7.5
VIEWPORT = {"width": 1920, "height": 1080}

# A headless Chromium runs in the same pod that serves live voting traffic.
# Capping concurrent exports per process to 1 (rather than letting requests
# queue silently for minutes, or stacking several browsers and starving that
# pod's CPU/memory during an actual event) keeps that footprint bounded and
# gives a caller an immediate, honest "try again shortly" instead of a hang.
_EXPORT_SLOT = asyncio.Semaphore(1)


async def _capture_steps(base_url: str, workflow_id: str, token: str, steps: list[WorkflowStep]) -> list[bytes]:
    if _EXPORT_SLOT.locked():
        raise HTTPException(503, "Another PPTX export is already running on this server. Please try again shortly.")
    async with _EXPORT_SLOT:
        return await _capture_steps_unguarded(base_url, workflow_id, token, steps)


async def _capture_steps_unguarded(base_url: str, workflow_id: str, token: str, steps: list[WorkflowStep]) -> list[bytes]:
    images: list[bytes] = []
    async with async_playwright() as playwright:
        browser = await playwright.chromium.launch(args=["--no-sandbox", "--disable-dev-shm-usage"])
        try:
            page = await browser.new_page(viewport=VIEWPORT)
            # WorkflowSceneRenderer.css fades every scene in (~0.6-0.7s,
            # staggered per element) and honors prefers-reduced-motion by
            # skipping that animation entirely -- request it here instead of
            # guessing a fixed wait, which under-shot on the hero/closing
            # scenes (their title+subtitle+pill chain hadn't finished fading
            # in yet) and produced washed-out, half-opacity slides.
            await page.emulate_media(reduced_motion="reduce")
            for step in steps:
                url = f"{base_url}/live/step-preview/{workflow_id}/{step.id}?token={token}"
                # networkidle is safe here (unlike the full TV display route)
                # because this preview page opens no realtime/SSE connection
                # -- nothing keeps the network permanently "busy". It's what
                # lets a video step's cross-origin YouTube iframe actually
                # finish loading its thumbnail before the screenshot, instead
                # of capturing a black box.
                await page.goto(url, wait_until="networkidle", timeout=30_000)
                await page.locator(".wf-scene").wait_for(timeout=20_000)
                # Small settle beat for web fonts / gradient paint, not animation.
                await page.wait_for_timeout(150)
                images.append(await page.screenshot())
        finally:
            await browser.close()
    return images


def _slide_notes(step: WorkflowStep, index: int, total: int) -> str:
    label = f'Slide {index + 1} of {total} · {step.step_type} · "{step.title}"'
    return f"{label}\n\n{step.presenter_notes}" if step.presenter_notes else label


def _assemble(steps: list[WorkflowStep], images: list[bytes]) -> bytes:
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
    blank_layout = presentation.slide_layouts[6]
    for index, (step, image_bytes) in enumerate(zip(steps, images)):
        slide = presentation.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            io.BytesIO(image_bytes), 0, 0,
            width=presentation.slide_width, height=presentation.slide_height,
        )
        slide.notes_slide.notes_text_frame.text = _slide_notes(step, index, len(steps))
    buffer = io.BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


async def render_workflow_pptx(base_url: str, workflow_id: str, token: str, steps: list[WorkflowStep]) -> bytes:
    images = await _capture_steps(base_url, workflow_id, token, steps)
    return _assemble(steps, images)
